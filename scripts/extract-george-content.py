"""
Extract text content from George's CFI Lesson Pack (.docx and .pptx files).
Saves structured JSON output per area for progress checkpointing.
"""

import json
import os
import sys
from pathlib import Path
from docx import Document
from pptx import Presentation

# Paths
BASE_PATH = Path(r"C:\Users\danrr\OneDrive - amazon.com\Desktop\CFI\_-_CFI_Lesson_Pack_from_George\Download All (PC) - CFI Lesson Pack")
OUTPUT_DIR = Path(r"C:\Users\danrr\OneDrive - amazon.com\Desktop\CFI\cfi-app-offline\scripts\extracted-content")

LESSONS_DIR = BASE_PATH / "1. CFI Lessons & PPTs" / "1. CFI Lessons" / "CFI Individual Files - Nov25"
EXTRA_DETAIL_DIR = BASE_PATH / "1. CFI Lessons & PPTs" / "3. CFI Extra Detail Lessons" / "CFI Extra Detail Individual Files - Nov25"
PPT_DIR = BASE_PATH / "1. CFI Lessons & PPTs" / "2. CFI PowerPoints" / "CFI Individual PPT Files - Nov25"

# Area folder name to Roman numeral mapping
AREA_MAP = {
    "1. Fundamentals of Instructing": "I",
    "2. Technical Subject Areas": "II",
    "3. Preflight Preparation": "III",
    "4. Preflight Lesson on a Maneuver": "IV",
    "5. Preflight Procedures": "V",
    "6. Airport Operations": "VI",
    "7. Takeoffs, Landings, & Go-Arounds": "VII",
    "8. Fundamentals of Flight": "VIII",
    "9. Maneuvers": "IX",
    "10. Slow Flight, Stalls, & Spins": "X",
    "11. Basic Instrument Maneuvers": "XI",
    "12. Emergency Operations": "XII",
    "14. Postflight Procedures": "XIV",
    "15. Appendix": "APP",
}


def extract_docx_content(filepath):
    """Extract all text content from a .docx file, preserving paragraph structure."""
    try:
        doc = Document(filepath)
        content = {
            "filename": filepath.name,
            "paragraphs": [],
            "tables": [],
            "headings": [],
            "full_text": "",
        }

        full_text_parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            para_data = {
                "text": text,
                "style": para.style.name if para.style else "Normal",
                "bold": any(run.bold for run in para.runs if run.bold),
            }
            content["paragraphs"].append(para_data)
            full_text_parts.append(text)

            # Track headings
            if para.style and "Heading" in para.style.name:
                content["headings"].append(text)

        content["full_text"] = "\n".join(full_text_parts)

        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                content["tables"].append(table_data)

        return content

    except Exception as e:
        return {"filename": filepath.name, "error": str(e), "full_text": ""}


def extract_pptx_content(filepath):
    """Extract text and speaker notes from a .pptx file."""
    try:
        prs = Presentation(filepath)
        content = {
            "filename": filepath.name,
            "slides": [],
            "full_text": "",
            "diagram_descriptions": [],
        }

        full_text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": i,
                "title": "",
                "text_content": [],
                "notes": "",
                "has_images": False,
                "has_shapes": False,
            }

            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text.strip()

            # Extract all text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_data["text_content"].append(text)
                            full_text_parts.append(text)

                # Check for images/diagrams
                if shape.shape_type == 13:  # Picture
                    slide_data["has_images"] = True
                elif hasattr(shape, "shapes"):  # Group shape
                    slide_data["has_shapes"] = True

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["notes"] = notes_text
                    full_text_parts.append(f"[NOTES] {notes_text}")

            # If slide has images, record as potential diagram
            if slide_data["has_images"] or slide_data["has_shapes"]:
                content["diagram_descriptions"].append({
                    "slide_number": i,
                    "title": slide_data["title"],
                    "context": " | ".join(slide_data["text_content"][:3]),
                })

            content["slides"].append(slide_data)

        content["full_text"] = "\n".join(full_text_parts)
        return content

    except Exception as e:
        return {"filename": filepath.name, "error": str(e), "full_text": ""}


def process_area(area_folder_name, area_roman):
    """Process all files for a single area and save checkpoint."""
    print(f"\n{'='*60}")
    print(f"Processing Area {area_roman}: {area_folder_name}")
    print(f"{'='*60}")

    area_data = {
        "area_number": area_roman,
        "area_name": area_folder_name,
        "standard_lessons": [],
        "extra_detail_lessons": [],
        "powerpoint_content": [],
    }

    # 1. Standard lessons
    standard_dir = LESSONS_DIR / area_folder_name
    if standard_dir.exists():
        docx_files = sorted(standard_dir.glob("*.docx"))
        print(f"  Standard lessons: {len(docx_files)} files")
        for f in docx_files:
            print(f"    Extracting: {f.name}")
            content = extract_docx_content(f)
            area_data["standard_lessons"].append(content)
    else:
        print(f"  Standard lessons: directory not found")

    # 2. Extra detail lessons
    extra_dir = EXTRA_DETAIL_DIR / area_folder_name
    if extra_dir.exists():
        docx_files = sorted(extra_dir.glob("*.docx"))
        print(f"  Extra detail lessons: {len(docx_files)} files")
        for f in docx_files:
            print(f"    Extracting: {f.name}")
            content = extract_docx_content(f)
            area_data["extra_detail_lessons"].append(content)
    else:
        print(f"  Extra detail lessons: directory not found")

    # 3. PowerPoints
    ppt_dir = PPT_DIR / area_folder_name
    if ppt_dir.exists():
        pptx_files = sorted(ppt_dir.glob("*.pptx"))
        print(f"  PowerPoints: {len(pptx_files)} files")
        for f in pptx_files:
            print(f"    Extracting: {f.name}")
            content = extract_pptx_content(f)
            area_data["powerpoint_content"].append(content)
    else:
        print(f"  PowerPoints: directory not found")

    # Save checkpoint
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_file = OUTPUT_DIR / f"area_{area_roman.replace(' ', '_')}.json"
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(area_data, f, ensure_ascii=False, indent=2)

    total = len(area_data["standard_lessons"]) + len(area_data["extra_detail_lessons"]) + len(area_data["powerpoint_content"])
    print(f"  ✓ Saved: {output_file.name} ({total} documents)")
    return area_data


def main():
    print("=" * 60)
    print("George's CFI Lesson Pack - Content Extraction")
    print("=" * 60)
    print(f"\nSource: {BASE_PATH}")
    print(f"Output: {OUTPUT_DIR}")
    print()

    # Verify paths exist
    if not BASE_PATH.exists():
        print(f"ERROR: Source path not found: {BASE_PATH}")
        sys.exit(1)

    all_areas = {}
    processed = 0
    errors = 0

    for area_folder, area_roman in AREA_MAP.items():
        try:
            area_data = process_area(area_folder, area_roman)
            all_areas[area_roman] = area_data
            processed += 1
        except Exception as e:
            print(f"  ERROR processing area {area_roman}: {e}")
            errors += 1

    # Save master index
    index_file = OUTPUT_DIR / "_extraction_index.json"
    index_data = {
        "total_areas_processed": processed,
        "errors": errors,
        "areas": {
            roman: {
                "standard_lessons": len(data["standard_lessons"]),
                "extra_detail_lessons": len(data["extra_detail_lessons"]),
                "powerpoints": len(data["powerpoint_content"]),
            }
            for roman, data in all_areas.items()
        },
    }
    with open(index_file, "w", encoding="utf-8") as f:
        json.dump(index_data, f, ensure_ascii=False, indent=2)

    print(f"\n{'='*60}")
    print(f"EXTRACTION COMPLETE")
    print(f"  Areas processed: {processed}")
    print(f"  Errors: {errors}")
    print(f"  Output: {OUTPUT_DIR}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
